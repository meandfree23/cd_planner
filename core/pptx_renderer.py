import re
import os
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------
# 디자인 토큰 (Design Tokens)
# ---------------------------------------------------------
COLOR_BG_DARK = RGBColor(10, 10, 10)
COLOR_TEXT_PRIMARY = RGBColor(255, 255, 255)
COLOR_TEXT_SECONDARY = RGBColor(180, 180, 180)
COLOR_TEXT_DARK = RGBColor(30, 30, 30)

COLOR_SCRIPT_BG = RGBColor(30, 30, 30)
COLOR_SCRIPT_TEXT = RGBColor(200, 200, 200)
COLOR_ACCENT = RGBColor(0, 122, 255) # Blue accent

FONT_FAMILY = 'Apple SD Gothic Neo'

# 12-Column Grid System (16:9 = 13.33 x 7.5 inches)
MARGIN_X = 1.0
MARGIN_Y = 0.8
CONTENT_WIDTH = 13.333 - (MARGIN_X * 2)
COL_WIDTH = CONTENT_WIDTH / 12

def get_x(col_start): return Inches(MARGIN_X + (col_start * COL_WIDTH))
def get_w(col_span): return Inches(col_span * COL_WIDTH)

def apply_text_style(tf, font_size, is_bold=False, color=COLOR_TEXT_PRIMARY, line_spacing=1.3, alignment=PP_ALIGN.LEFT):
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        paragraph.line_spacing = line_spacing
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = font_size
            run.font.bold = is_bold
            run.font.color.rgb = color

def create_pptx_file(ppt_content, filepath, brief_input):
    from core.nodes import get_openai_llm
    from openai import OpenAI
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 1. DALL-E 배경 생성 (공통 텍스처)
    cover_image_stream = None
    try:
        client = OpenAI()
        llm = get_openai_llm()
        prompt_resp = llm.invoke(f"다음 기획서 브리프를 읽고, 시각적 영문 키워드(단어 2~3개) 추출해줘. \n{brief_input[:500]}")
        keyword = prompt_resp.content.strip()
        
        # 다크모드 미니멀 텍스처 배경
        dalle_prompt = f"A very subtle, extremely dark, minimalist abstract background texture for {keyword}. Dark mode, corporate luxury, no text, no sharp objects, smooth gradients --ar 16:9"
        
        print("🎨 Generating DALL-E 3 background...")
        response = client.images.generate(
            model="dall-e-3",
            prompt=dalle_prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        img_response = requests.get(response.data[0].url)
        cover_image_stream = img_response.content
    except Exception as e:
        print(f"DALL-E generation failed: {e}")

    blocks = re.split(r'(?i)(?:##\s*)?slide\s*\d+.*?\n', ppt_content)
    if len(blocks) <= 1:
        blocks = ppt_content.split('\n\n\n')
        
    for i, block in enumerate(blocks):
        block = block.strip()
        if not block: continue
        
        lines = block.split('\n')
        title_text = lines[0].replace('#', '').strip()[:100] if lines else "Slide"
        
        # Semantic Parser
        copy_text = []
        body_text = []
        script_text = []
        
        current_context = "body"
        
        for line in lines[1:]:
            line_clean = line.strip()
            if not line_clean:
                continue
                
            if "[카피]" in line_clean or "[핵심]" in line_clean:
                current_context = "copy"
                copy_text.append(line_clean.replace("[카피]", "").replace("[핵심]", "").replace("-", "").strip())
            elif "[PT 스크립트]" in line_clean or "[스크립트]" in line_clean:
                current_context = "script"
                script_text.append(line_clean.replace("[PT 스크립트]", "").replace("[스크립트]", "").replace("-", "").strip())
            elif "[설명]" in line_clean:
                current_context = "body"
            else:
                if current_context == "copy": copy_text.append(line_clean)
                elif current_context == "script": script_text.append(line_clean)
                else: body_text.append(line_clean)
                
        # 슬라이드 생성 (항상 Blank 레이아웃 사용)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # DALL-E 배경 깔기 (맨 밑 레이어)
        if cover_image_stream:
            pic = slide.shapes.add_picture(BytesIO(cover_image_stream), 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            # Fallback 다크모드 배경
            bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height) # 1=Rectangle
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_BG_DARK
            bg.line.fill.background()
            
        # 표지 슬라이드 (첫 슬라이드)
        if i == 0 or (i == 1 and not blocks[0].strip()):
            txBox = slide.shapes.add_textbox(get_x(0), Inches(3.0), get_w(10), Inches(2))
            tf = txBox.text_frame
            tf.text = title_text
            apply_text_style(tf, Pt(48), is_bold=True, color=COLOR_TEXT_PRIMARY)
            
            subBox = slide.shapes.add_textbox(get_x(0), Inches(5.0), get_w(8), Inches(1))
            stf = subBox.text_frame
            stf.text = "AI Creative Director Planner"
            apply_text_style(stf, Pt(18), color=COLOR_ACCENT)
            continue
            
        # 일반 슬라이드 - 헤더 (좌상단 1~12열)
        titleBox = slide.shapes.add_textbox(get_x(0), Inches(0.6), get_w(12), Inches(0.8))
        ttf = titleBox.text_frame
        ttf.text = title_text
        apply_text_style(ttf, Pt(24), is_bold=True, color=COLOR_ACCENT)
        
        # 카피가 있을 경우 (히어로 섹션, 1~6열)
        if copy_text:
            copyBox = slide.shapes.add_textbox(get_x(0), Inches(1.8), get_w(6), Inches(2.5))
            copyBox.text_frame.text = "\n".join(copy_text)
            apply_text_style(copyBox.text_frame, Pt(32), is_bold=True, line_spacing=1.4)
            
            # 본문은 우측 (7~12열)
            if body_text:
                bodyBox = slide.shapes.add_textbox(get_x(6.5), Inches(1.8), get_w(5.5), Inches(3))
                bodyBox.text_frame.text = "\n".join(body_text)
                apply_text_style(bodyBox.text_frame, Pt(16), color=COLOR_TEXT_SECONDARY, line_spacing=1.5)
        else:
            # 카피가 없으면 본문을 넓게 (1~9열)
            if body_text:
                bodyBox = slide.shapes.add_textbox(get_x(0), Inches(1.8), get_w(9), Inches(3.5))
                bodyBox.text_frame.text = "\n".join(body_text)
                apply_text_style(bodyBox.text_frame, Pt(18), color=COLOR_TEXT_PRIMARY, line_spacing=1.5)
                
        # 스크립트 박스 (하단 1~12열, 디자인된 박스 형태)
        if script_text:
            script_y = Inches(5.8)
            script_h = Inches(1.2)
            
            # 스크립트 배경 도형
            s_bg = slide.shapes.add_shape(1, get_x(0), script_y, get_w(12), script_h)
            s_bg.fill.solid()
            s_bg.fill.fore_color.rgb = COLOR_SCRIPT_BG
            s_bg.line.fill.background()
            
            sBox = slide.shapes.add_textbox(get_x(0) + Inches(0.2), script_y + Inches(0.1), get_w(11.6), script_h - Inches(0.2))
            sBox.text_frame.text = "💡 발표자 노트: " + " ".join(script_text)
            apply_text_style(sBox.text_frame, Pt(12), color=COLOR_SCRIPT_TEXT, line_spacing=1.2)
            
        # 페이지 번호
        numBox = slide.shapes.add_textbox(get_x(11.5), Inches(6.8), get_w(0.5), Inches(0.3))
        numBox.text_frame.text = str(i+1)
        apply_text_style(numBox.text_frame, Pt(11), color=COLOR_TEXT_SECONDARY, alignment=PP_ALIGN.RIGHT)

    prs.save(filepath)
    return True
